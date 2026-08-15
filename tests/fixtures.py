"""テスト用の Office ファイルを組み立てるヘルパ。

v0.1 では手書きの XML を使っていたが、変換の本体が markitdown になったことで
**本物のライブラリが読める妥当なパッケージ**でないとテストできなくなった
（[Content_Types].xml などが無いと開けない）。

そこで python-docx / openpyxl / python-pptx に器を作らせ、
検査対象にしたい XML（グラフ・テキストボックス等）だけを後から差し込む。
差し込む XML はテストに直接書くので、「どの構造を検出したいのか」は読めば分かる。
"""

from __future__ import annotations

import io
import shutil
import struct
import zipfile
import zlib
from pathlib import Path

A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def png(color: tuple[int, int, int] = (0, 120, 200)) -> bytes:
    """最小の PNG。python-docx が画像として受け付ける必要があるので本物を作る。"""

    def chunk(tag: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(tag + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

    width = height = 4
    header = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + bytes(color) * width for _ in range(height))
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", header)
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


def para(text: str, *, style: str | None = None) -> str:
    """段落の XML。style は python-docx が用意する組み込みスタイル ID。"""
    ppr = f'<w:pPr><w:pStyle w:val="{style}"/></w:pPr>' if style else ""
    return f'<w:p xmlns:w="{W}">{ppr}<w:r><w:t xml:space="preserve">{text}</w:t></w:r></w:p>'


def docx(
    path: Path,
    body: str = "",
    *,
    rels: dict[str, str] | None = None,
    media: dict[str, bytes] | None = None,
    title: str | None = None,
    picture: bytes | None = None,
) -> Path:
    """python-docx で作った文書に、任意の本文 XML を差し込む。

    picture を渡すと本物の画像を 1 枚入れる（画像の取り出しを試すため）。
    """
    import docx as python_docx

    document = python_docx.Document()
    if title:
        document.core_properties.title = title
    if picture is not None:
        document.add_picture(io.BytesIO(picture))
    buffer = io.BytesIO()
    document.save(buffer)

    parts = _read_zip(buffer.getvalue())
    parts["word/document.xml"] = _insert_into_body(parts["word/document.xml"], body)
    if rels:
        parts["word/_rels/document.xml.rels"] = _add_relationships(
            parts["word/_rels/document.xml.rels"], rels
        )
    for name, data in (media or {}).items():
        parts[f"word/media/{name}"] = data
    return _write_zip(path, parts)


def xlsx(
    path: Path, sheets: dict[str, list[list[object]]], *, hidden: set[str] | None = None
) -> Path:
    """openpyxl でブックを作る。hidden に入れたシートは非表示にする。"""
    from openpyxl import Workbook

    book = Workbook()
    book.remove(book.active)
    for name, grid in sheets.items():
        sheet = book.create_sheet(name)
        for row in grid:
            sheet.append(list(row))
        if hidden and name in hidden:
            sheet.sheet_state = "hidden"
    book.save(path)
    return path


def pptx(path: Path, slides: list[dict]) -> Path:
    """python-pptx でスライドを作る。slides は {"title":…, "bullets":[(level, text)]}。"""
    from pptx import Presentation

    deck = Presentation()
    for slide in slides:
        added = deck.slides.add_slide(deck.slide_layouts[1])
        added.shapes.title.text = slide.get("title", "")
        bullets = slide.get("bullets") or []
        if bullets:
            frame = added.placeholders[1].text_frame
            frame.text = bullets[0][1]
            for level, text in bullets[1:]:
                paragraph = frame.add_paragraph()
                paragraph.text = text
                paragraph.level = level
        if slide.get("notes"):
            added.notes_slide.notes_text_frame.text = slide["notes"]
    deck.save(path)
    return path


# --------------------------------------------------------------------------
# ZIP の組み替え
# --------------------------------------------------------------------------


def _read_zip(data: bytes) -> dict[str, bytes]:
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        return {name: zf.read(name) for name in zf.namelist()}


def _write_zip(path: Path, parts: dict[str, bytes]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, data in parts.items():
            zf.writestr(name, data)
    shutil.move(tmp, path)
    return path


def _insert_into_body(document_xml: bytes, body: str) -> bytes:
    """`<w:body>` の直後に XML を差し込む（本文の先頭に置く）。"""
    if not body:
        return document_xml
    text = document_xml.decode("utf-8")
    marker = "<w:body>"
    index = text.index(marker) + len(marker)
    return (text[:index] + body + text[index:]).encode("utf-8")


def _add_relationships(rels_xml: bytes, mapping: dict[str, str]) -> bytes:
    text = rels_xml.decode("utf-8")
    items = []
    for rid, target in mapping.items():
        mode = ' TargetMode="External"' if "://" in target else ""
        items.append(
            f'<Relationship Id="{rid}" Target="{target}" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"'
            f"{mode}/>"
        )
    items = "".join(items)
    return text.replace("</Relationships>", items + "</Relationships>").encode("utf-8")
