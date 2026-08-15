"""回帰テスト用の実ファイル（.docx / .xlsx / .pptx）を生成する。

手書きの XML（`tests/fixtures.py`）では、Word/Excel/PowerPoint が実際に吐く
スタイル定義・番号定義・関係ファイルの作られ方までは再現できない。
そこで python-docx / openpyxl / python-pptx に本物のファイルを作らせ、
その出力を固定して回帰テストの土台にする（タスク T-01）。

    .venv/bin/python tools/build_fixtures.py

生成物は `tests/assets/` に置き、リポジトリにコミットする。
このスクリプトは開発時にしか使わないので、mdconv 本体はこれらのライブラリに依存しない。
"""

from __future__ import annotations

import struct
import zlib
from pathlib import Path

ASSETS = Path(__file__).resolve().parent.parent / "tests" / "assets"


def png(width: int = 4, height: int = 4, color: tuple[int, int, int] = (0, 120, 200)) -> bytes:
    """依存なしで最小の PNG を作る（Pillow に頼らず内容を固定するため）。"""

    def chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    header = struct.pack(">IIBBBBB", width, height, 8, 2, 0, 0, 0)
    raw = b"".join(b"\x00" + bytes(color) * width for _ in range(height))
    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", header)
        + chunk(b"IDAT", zlib.compress(raw))
        + chunk(b"IEND", b"")
    )


def build_docx(path: Path) -> None:
    import docx

    doc = docx.Document()
    doc.core_properties.title = "変換テスト用ドキュメント"

    doc.add_heading("プロジェクト提案書", level=0)  # Title スタイル
    doc.add_heading("背景", level=1)

    p = doc.add_paragraph("社内資料は ")
    p.add_run("Office 形式").bold = True
    p.add_run(" と PDF に閉じ込められている。")
    p.add_run("検索できない").italic = True
    p.add_run("のが最大の問題。")

    doc.add_heading("課題", level=2)
    doc.add_paragraph("資料が探せない", style="List Bullet")
    doc.add_paragraph("差分が見えない", style="List Bullet 2")
    doc.add_paragraph("AI に渡しにくい", style="List Bullet")

    doc.add_heading("進め方", level=2)
    doc.add_paragraph("現状を調べる", style="List Number")
    doc.add_paragraph("変換ツールを作る", style="List Number")

    doc.add_paragraph("引用のような一文。", style="Quote")

    table = doc.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    for row, values in enumerate(
        [["形式", "対応", "備考"], ["Word", "済", "見出し・表"], ["PDF", "一部", "表が苦手"]]
    ):
        for col, value in enumerate(values):
            table.cell(row, col).text = value

    doc.add_heading("参考", level=1)
    doc.add_paragraph("詳しくは社内 Wiki を参照。")

    image = ASSETS / "_tmp_image.png"
    image.write_bytes(png())
    doc.add_picture(str(image))
    image.unlink()

    doc.save(path)


def build_xlsx(path: Path) -> None:
    from openpyxl import Workbook

    book = Workbook()

    sales = book.active
    sales.title = "売上"
    sales.append(["月", "金額", "前月比"])
    sales.append(["1月", 1200, None])
    sales.append(["2月", 1500, "=B3/B2"])
    sales.append(["3月", 900, "=B4/B3"])

    members = book.create_sheet("名簿")
    members.append(["田中", "営業", 30])
    members.append(["鈴木", "開発", 25])

    empty = book.create_sheet("空シート")
    empty.sheet_state = "visible"

    secret = book.create_sheet("内部用")
    secret.append(["原価", 800])
    secret.sheet_state = "hidden"

    book.save(path)
    # openpyxl は数式しか書けないが、実際の .xlsx には Excel が計算した値も入っている。
    # mdconv はその値を出力する仕様（FR-102）なので、本物に合わせて補う
    inject_cached_values(path, "xl/worksheets/sheet1.xml", {"C3": "1.25", "C4": "0.6"})


def inject_cached_values(path: Path, part: str, values: dict[str, str]) -> None:
    """数式セルに計算結果 <v> を差し込む（Excel が保存した状態を再現する）。"""
    import re
    import shutil
    import zipfile

    with zipfile.ZipFile(path) as src:
        parts = {name: src.read(name) for name in src.namelist()}

    xml = parts[part].decode("utf-8")
    for ref, value in values.items():
        # openpyxl は計算結果を持たないので <v></v> を空のまま書く。そこを埋める
        pattern = re.compile(rf'(<c r="{ref}"[^>]*>)(<f>[^<]*</f>)(<v\s*/>|<v></v>)?')
        xml, count = pattern.subn(rf"\1\2<v>{value}</v>", xml)
        if count != 1:
            raise RuntimeError(f"{ref} の数式セルが見つかりません（{count} 件）")
    parts[part] = xml.encode("utf-8")

    tmp = path.with_suffix(".tmp")
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as out:
        for name, data in parts.items():
            out.writestr(name, data)
    shutil.move(tmp, path)


def build_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    deck = Presentation()

    slide = deck.slides.add_slide(deck.slide_layouts[1])
    slide.shapes.title.text = "はじめに"
    body = slide.placeholders[1].text_frame
    body.text = "背景"
    for text, level in (("資料が探せない", 1), ("目的", 0), ("Markdown 化する", 1)):
        para = body.add_paragraph()
        para.text = text
        para.level = level
    slide.notes_slide.notes_text_frame.text = "ここで事例を 1 つ話す"

    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "対応状況"
    table = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    for row, values in enumerate([["形式", "状態"], ["Word", "済"], ["PDF", "一部"]]):
        for col, value in enumerate(values):
            table.cell(row, col).text = value

    slide = deck.slides.add_slide(deck.slide_layouts[6])  # 白紙（タイトルなし）
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.text_frame.text = "ご清聴ありがとうございました"

    deck.save(path)


def main() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    build_docx(ASSETS / "sample.docx")
    build_xlsx(ASSETS / "sample.xlsx")
    build_pptx(ASSETS / "sample.pptx")
    print(f"生成しました: {', '.join(sorted(p.name for p in ASSETS.glob('sample.*')))}")


if __name__ == "__main__":
    main()
